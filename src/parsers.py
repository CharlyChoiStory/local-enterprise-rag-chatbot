from __future__ import annotations

"""
로컬 문서 파서.
parse_document(path) 하나로 PDF/TXT/DOCX/PPTX를 처리하고
표준 구조를 반환한다.  파싱 실패 시 앱 전체가 죽지 않고 error 정보를 반환한다.
"""

import hashlib
import logging
from pathlib import Path
from typing import Any

logger = logging.getLogger(__name__)

# 카테고리는 상위 폴더명에서 자동 감지
_KNOWN_CATEGORIES = {"hr", "finance", "welfare", "it", "policy"}

# 최소 의미 있는 텍스트 길이 (글자)
_MIN_SECTION_LEN = 10


# ── 내부 유틸 ────────────────────────────────────────────────

def _file_hash(path: Path) -> str:
    h = hashlib.sha256()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(65536), b""):
            h.update(chunk)
    return h.hexdigest()[:16]


def _detect_category(path: Path) -> str:
    for part in path.parts:
        if part.lower() in _KNOWN_CATEGORIES:
            return part.lower()
    return "general"


def _clean(text: str) -> str:
    """연속 공백·빈 줄 정리."""
    import re
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


# ── 포맷별 파서 ──────────────────────────────────────────────

def _parse_pdf(path: Path) -> list[dict]:
    import fitz  # PyMuPDF

    sections = []
    doc = fitz.open(str(path))
    for page in doc:
        text = _clean(page.get_text())
        if len(text) >= _MIN_SECTION_LEN:
            sections.append({"location": f"page {page.number + 1}", "text": text})
    doc.close()
    return sections


def _parse_txt(path: Path) -> list[dict]:
    from charset_normalizer import from_path

    result = from_path(str(path)).best()
    if result is None:
        raise ValueError("인코딩 감지 실패")
    text = _clean(str(result))
    if len(text) < _MIN_SECTION_LEN:
        return []
    import re
    paragraphs = [p.strip() for p in re.split(r"\n{2,}", text) if len(p.strip()) >= _MIN_SECTION_LEN]
    return [{"location": f"paragraph {i + 1}", "text": p} for i, p in enumerate(paragraphs)]


def _parse_docx(path: Path) -> list[dict]:
    from docx import Document

    doc = Document(str(path))
    sections = []
    para_num = 0
    for para in doc.paragraphs:
        text = _clean(para.text)
        if len(text) >= _MIN_SECTION_LEN:
            para_num += 1
            sections.append({"location": f"paragraph {para_num}", "text": text})
    return sections


def _parse_pptx(path: Path) -> list[dict]:
    from pptx import Presentation

    prs = Presentation(str(path))
    sections = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        parts.append(t)
        text = _clean("\n".join(parts))
        if len(text) >= _MIN_SECTION_LEN:
            sections.append({"location": f"slide {slide_num}", "text": text})
    return sections


# ── 공개 API ─────────────────────────────────────────────────

def parse_document(path: str | Path) -> dict[str, Any]:
    """
    단일 파일을 파싱해 표준 구조로 반환한다.

    반환값:
    {
        "document_id": str,          # SHA-256 앞 16자
        "filename":    str,
        "file_path":   str,
        "file_type":   "pdf"|"txt"|"docx"|"pptx",
        "category":    str,
        "sections":    [{"location": str, "text": str}, ...],
        "error":       str | None    # 파싱 실패 시 메시지
    }
    """
    path = Path(path)
    base: dict[str, Any] = {
        "document_id": "",
        "filename":    path.name,
        "file_path":   str(path),
        "file_type":   path.suffix.lstrip(".").lower(),
        "category":    _detect_category(path),
        "sections":    [],
        "error":       None,
    }

    try:
        base["document_id"] = _file_hash(path)
        ext = path.suffix.lower()

        if ext == ".pdf":
            base["sections"] = _parse_pdf(path)
        elif ext == ".txt":
            base["sections"] = _parse_txt(path)
        elif ext == ".docx":
            base["sections"] = _parse_docx(path)
        elif ext == ".pptx":
            base["sections"] = _parse_pptx(path)
        else:
            base["error"] = f"지원하지 않는 파일 형식: {ext}"

    except Exception as exc:
        logger.warning("파싱 실패 [%s]: %s", path.name, exc)
        base["error"] = str(exc)

    return base


def scan_documents(docs_dir: str | Path) -> list[dict[str, Any]]:
    """
    docs_dir 하위의 PDF/TXT/DOCX/PPTX 파일을 모두 파싱해 반환한다.
    실패한 파일도 error 필드에 기록한 채 포함된다.
    """
    docs_dir = Path(docs_dir)
    results = []
    for ext in ("*.pdf", "*.txt", "*.docx", "*.pptx"):
        for path in sorted(docs_dir.rglob(ext)):
            results.append(parse_document(path))
    return results
